import os
import re
import io
import copy  # For deep copying shapes
import matplotlib.pyplot as plt
from PIL import Image
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor  # Import RGBColor for font color
from google.auth.transport.requests import Request  # For refreshing credentials
import tkinter as tk
from tkinter import filedialog, messagebox, ttk, colorchooser, font
import threading

# Google Sheets API scope
SCOPES = ['https://www.googleapis.com/auth/spreadsheets.readonly']


# GUI Class
class PowerPointGeneratorGUI:
    def __init__(self, root):
        self.root = root
        self.root.title("PowerPoint Generator")
        self.create_widgets()


    def create_widgets(self):
        padding = {'padx': 10, 'pady': 5}

        # Create a notebook for organizing widgets into tabs
        notebook = ttk.Notebook(self.root)
        notebook.grid(row=0, column=0, sticky='nsew', padx=10, pady=10)
        self.root.grid_rowconfigure(0, weight=1)
        self.root.grid_columnconfigure(0, weight=1)

        # Tab 1: File Selection
        file_frame = ttk.Frame(notebook)
        notebook.add(file_frame, text='File Selection')

        # Credentials File
        self.credentials_path = tk.StringVar()
        ttk.Label(file_frame, text="Credentials JSON:").grid(row=0, column=0, sticky='e', **padding)
        self.credentials_path_entry = ttk.Entry(file_frame, textvariable=self.credentials_path, width=50)
        self.credentials_path_entry.grid(row=0, column=1, **padding)
        self.browse_credentials_button = ttk.Button(file_frame, text="Browse", command=self.browse_credentials)
        self.browse_credentials_button.grid(row=0, column=2, **padding)

        # Token File
        self.token_path = tk.StringVar(value='token.json')  # Default value
        ttk.Label(file_frame, text="Token JSON:").grid(row=1, column=0, sticky='e', **padding)
        self.token_path_entry = ttk.Entry(file_frame, textvariable=self.token_path, width=50)
        self.token_path_entry.grid(row=1, column=1, **padding)
        self.browse_token_button = ttk.Button(file_frame, text="Browse", command=self.browse_token)
        self.browse_token_button.grid(row=1, column=2, **padding)

        # Template PowerPoint
        self.template_pptx = tk.StringVar()
        ttk.Label(file_frame, text="Template PPTX:").grid(row=2, column=0, sticky='e', **padding)
        self.template_pptx_entry = ttk.Entry(file_frame, textvariable=self.template_pptx, width=50)
        self.template_pptx_entry.grid(row=2, column=1, **padding)
        self.browse_template_button = ttk.Button(file_frame, text="Browse", command=self.browse_template)
        self.browse_template_button.grid(row=2, column=2, **padding)

        # Spreadsheet ID
        self.spreadsheet_id = tk.StringVar()
        ttk.Label(file_frame, text="Spreadsheet ID:").grid(row=3, column=0, sticky='e', **padding)
        self.spreadsheet_id_entry = ttk.Entry(file_frame, textvariable=self.spreadsheet_id, width=50)
        self.spreadsheet_id_entry.grid(row=3, column=1, columnspan=2, **padding)

        # Output PowerPoint
        self.output_pptx = tk.StringVar(value='output.pptx')  # Default value
        ttk.Label(file_frame, text="Output PPTX:").grid(row=4, column=0, sticky='e', **padding)
        self.output_pptx_entry = ttk.Entry(file_frame, textvariable=self.output_pptx, width=50)
        self.output_pptx_entry.grid(row=4, column=1, **padding)
        self.browse_output_button = ttk.Button(file_frame, text="Browse", command=self.browse_output)
        self.browse_output_button.grid(row=4, column=2, **padding)

        # Tab 2: Font Settings
        font_frame = ttk.Frame(notebook)
        notebook.add(font_frame, text='Font Settings')

        # Retrieve all available fonts
        available_fonts = sorted(font.families())

        # Define sections for each text element
        self.font_settings = {
            'Question': {},
            'Option': {},
            'Board/Institute': {},
            'Answer': {}
        }

        # For each text element, create a frame with font options
        for idx, element in enumerate(self.font_settings.keys()):
            section = ttk.LabelFrame(font_frame, text=f"{element} Font")
            section.grid(row=idx, column=0, sticky='ew', padx=10, pady=10, ipadx=10, ipady=10)

            # Font Name
            ttk.Label(section, text="Font Name:").grid(row=0, column=0, sticky='e', padx=5, pady=5)
            font_name = tk.StringVar(value='Arial')
            font_dropdown = ttk.Combobox(section, textvariable=font_name, values=available_fonts, state='readonly')
            font_dropdown.grid(row=0, column=1, padx=5, pady=5)
            self.font_settings[element]['name'] = font_name

            # Font Size
            ttk.Label(section, text="Font Size:").grid(row=1, column=0, sticky='e', padx=5, pady=5)
            default_size = 18 if element == 'Question' else 14
            font_size = tk.IntVar(value=default_size)
            font_size_spin = ttk.Spinbox(section, from_=8, to=72, textvariable=font_size, width=5)
            font_size_spin.grid(row=1, column=1, sticky='w', padx=5, pady=5)
            self.font_settings[element]['size'] = font_size

            # Bold
            bold_var = tk.BooleanVar(value=True if element == 'Question' else False)
            bold_check = ttk.Checkbutton(section, text="Bold", variable=bold_var)
            bold_check.grid(row=2, column=0, sticky='e', padx=5, pady=5)
            self.font_settings[element]['bold'] = bold_var

            # Italic
            italic_var = tk.BooleanVar(value=False)
            italic_check = ttk.Checkbutton(section, text="Italic", variable=italic_var)
            italic_check.grid(row=2, column=1, sticky='w', padx=5, pady=5)
            self.font_settings[element]['italic'] = italic_var

            # Font Color
            ttk.Label(section, text="Font Color:").grid(row=3, column=0, sticky='e', padx=5, pady=5)
            color_var = tk.StringVar(value='#000000')  # Default black
            color_button = ttk.Button(section, text="Choose Color", command=lambda var=color_var, sec=section: self.choose_color(var, sec))
            color_button.grid(row=3, column=1, sticky='w', padx=5, pady=5)
            self.font_settings[element]['color'] = color_var

        # Tab 1 continued: Run Button and Status Display
        # Place Run Button and Status Display outside the notebook for better visibility
        self.generate_button = ttk.Button(file_frame, text="Generate PowerPoint", command=self.start_generation_thread)
        self.generate_button.grid(row=5, column=0, columnspan=3, pady=20)

        # Status Display
        self.status_text = tk.Text(self.root, height=15, width=100, state='disabled')
        self.status_text.grid(row=1, column=0, padx=10, pady=10)

        # Maintain a list of target widgets for selective disabling/enabling
        self.target_widgets = [
            self.credentials_path_entry,
            self.token_path_entry,
            self.template_pptx_entry,
            self.spreadsheet_id_entry,
            self.output_pptx_entry,
            self.browse_credentials_button,
            self.browse_token_button,
            self.browse_template_button,
            self.browse_output_button,
            self.generate_button
        ]

    def choose_color(self, color_var, section):
        color_code = colorchooser.askcolor(title="Choose font color")
        if color_code:
            color_var.set(color_code[1])  # Hex color code
            # Optionally, you can update a label or indicator to show the selected color

    def browse_credentials(self):
        file_path = filedialog.askopenfilename(title="Select credentials.json", filetypes=[("JSON Files", "*.json")])
        if file_path:
            self.credentials_path.set(file_path)

    def browse_token(self):
        file_path = filedialog.askopenfilename(title="Select token.json", filetypes=[("JSON Files", "*.json")])
        if file_path:
            self.token_path.set(file_path)

    def browse_template(self):
        file_path = filedialog.askopenfilename(title="Select Template PPTX", filetypes=[("PowerPoint Files", "*.pptx")])
        if file_path:
            self.template_pptx.set(file_path)

    def browse_output(self):
        file_path = filedialog.asksaveasfilename(title="Save Output PPTX", defaultextension=".pptx",
                                                 filetypes=[("PowerPoint Files", "*.pptx")])
        if file_path:
            self.output_pptx.set(file_path)

    def log_status(self, message):
        self.status_text.configure(state='normal')
        self.status_text.insert(tk.END, message + "\n")
        self.status_text.configure(state='disabled')
        self.status_text.see(tk.END)

    def start_generation_thread(self):
        # Start the generation in a separate thread to keep GUI responsive
        generation_thread = threading.Thread(target=self.run_generation)
        generation_thread.start()

    def run_generation(self):
        # Disable target widgets to prevent multiple clicks
        self.set_widgets_state('disabled')
        self.log_status("Starting PowerPoint generation...")

        try:
            # Retrieve inputs from GUI
            credentials_path = self.credentials_path.get()
            token_path = self.token_path.get()
            template_pptx = self.template_pptx.get()
            output_pptx = self.output_pptx.get()
            spreadsheet_id_input = self.spreadsheet_id.get()

            # Validate required inputs
            if not all([credentials_path, template_pptx, spreadsheet_id_input]):
                raise ValueError("Please provide all required inputs: Credentials JSON, Template PPTX, and Spreadsheet ID.")

            # Collect font settings
            font_settings_collected = {}
            for element, settings in self.font_settings.items():
                font_settings_collected[element] = {
                    'name': settings['name'].get(),
                    'size': settings['size'].get(),
                    'bold': settings['bold'].get(),
                    'italic': settings['italic'].get(),
                    'color': settings['color'].get()
                }

            # Call the main function with updated parameters
            main_gui_wrapper(
                credentials_path=credentials_path,
                token_path=token_path,
                spreadsheet_id_input=spreadsheet_id_input,
                template_pptx=template_pptx,
                output_pptx=output_pptx,
                font_settings=font_settings_collected,
                log_func=self.log_status
            )

        except Exception as e:
            messagebox.showerror("Error", str(e))
            self.log_status(f"Error: {e}")
        finally:
            # Re-enable the widgets
            self.set_widgets_state('normal')

    def set_widgets_state(self, state):
        for widget in self.target_widgets:
            try:
                widget.configure(state=state)
            except tk.TclError:
                # Widget does not support 'state'; skip it
                pass

    def close(self):
        self.root.quit()


def authenticate_gui(credentials_path_param, token_path_param, log_func=None):
    creds = None
    # Token file stores the user's access and refresh tokens
    if os.path.exists(token_path_param):
        creds = Credentials.from_authorized_user_file(token_path_param, SCOPES)
        if log_func:
            log_func(f"Loaded credentials from {token_path_param}")
    # If credentials are not available or invalid, initiate the OAuth flow
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
            if log_func:
                log_func("Credentials refreshed.")
        else:
            flow = InstalledAppFlow.from_client_secrets_file(credentials_path_param, SCOPES)
            creds = flow.run_local_server(port=0)
            if log_func:
                log_func("Obtained new credentials via OAuth flow.")
        # Save the credentials for future runs
        with open(token_path_param, 'w') as token:
            token.write(creds.to_json())
            if log_func:
                log_func(f"Saved credentials to {token_path_param}")
    return creds


def main_gui_wrapper(credentials_path, token_path, spreadsheet_id_input, template_pptx, output_pptx, font_settings, log_func=None):
    creds = authenticate_gui(credentials_path, token_path, log_func)
    service = build('sheets', 'v4', credentials=creds)
    log_func("Authenticated with Google Sheets API.")

    range_name = 'Sheet1!A1:L'  # Adjusted to cover columns H, I, K, and L
    log_func(f"Fetching data from Spreadsheet ID: {spreadsheet_id_input}")

    # Read data from Google Sheets
    sheet = service.spreadsheets()
    try:
        result = sheet.values().get(spreadsheetId=spreadsheet_id_input, range=range_name).execute()
    except Exception as e:
        log_func(f"API Error: {e}")
        raise e
    values = result.get('values', [])

    log_func(f"Raw data fetched: {values}")  # Log raw data for debugging

    if not values:
        log_func('No data found in the Google Sheet.')
        return

    log_func(f"Retrieved {len(values)-1} rows of data.")

    # Extract headers and data
    headers = values[0]
    data_rows = values[1:]

    # Load the PowerPoint template
    if os.path.exists(template_pptx):
        prs = Presentation(template_pptx)
        log_func(f"Loaded template PowerPoint from {template_pptx}")
    else:
        raise FileNotFoundError(f"Template file {template_pptx} not found.")

    # List available slide layouts
    log_func("Available slide layouts:")
    for index, layout in enumerate(prs.slide_layouts):
        log_func(f"Index {index}: {layout.name}")

    # Assuming the first slide in the template is the one with placeholders
    template_slide_index = 0  # Adjust if necessary

    # Modify to target the correct columns for options (H, I, K, L)
    option_columns = ['Option_ক', 'Option_খ', 'Option_গ', 'Option_ঘ']
    log_func(f"Identified option columns: {option_columns}")

    for row_num, row in enumerate(data_rows, start=1):
        row_data = dict(zip(headers, row))
        question_preview = row_data.get('Question', '')[:30] + "..." if len(row_data.get('Question', '')) > 30 else row_data.get('Question', '')
        log_func(f"Processing row {row_num}: {question_preview}")

        # Duplicate the template slide
        slide = duplicate_slide(prs, template_slide_index)
        log_func("Duplicated template slide.")

        # Insert content into text boxes based on placeholder text
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if '[question]' in text:
                    question_text = row_data.get('Question', '')
                    question_text_no_latex = remove_latex(question_text)
                    shape.text_frame.text = f"{row_data.get('Serial Number', '')}. {question_text_no_latex}"

                    # Apply font settings for Question
                    apply_font_settings(shape, font_settings['Question'], log_func=log_func)

                    # Handle equations in question
                    question_latex_codes = extract_latex(question_text)
                    left = shape.left
                    top = shape.top + shape.height + Inches(0.1)
                    for code in question_latex_codes:
                        image = latex_to_image(code)
                        image_stream = io.BytesIO()
                        image.save(image_stream, format='PNG')
                        image_stream.seek(0)

                        # Add picture with desired size
                        pic_width = Inches(4)  # Adjust as needed
                        pic_height = Inches(1)  # Adjust as needed
                        pic = slide.shapes.add_picture(image_stream, left, top, width=pic_width, height=pic_height)
                        top += pic.height + Inches(0.1)
                        log_func(f"Added LaTeX image for question: {code}")
                elif '[option_ক]' in text:
                    option_text = row_data.get('Option_ক', '')
                    option_text_no_latex = remove_latex(option_text)
                    shape.text_frame.text = f"Option ক. {option_text_no_latex}"

                    # Apply font settings for Option_ক
                    apply_font_settings(shape, font_settings['Option'], log_func=log_func)

                    # Handle equations in option
                    option_latex_codes = extract_latex(option_text)
                    left = shape.left
                    top = shape.top + shape.height + Inches(0.1)
                    for code in option_latex_codes:
                        image = latex_to_image(code)
                        image_stream = io.BytesIO()
                        image.save(image_stream, format='PNG')
                        image_stream.seek(0)

                        # Add picture with desired size
                        pic_width = Inches(3)  # Adjust as needed
                        pic_height = Inches(0.8)  # Adjust as needed
                        pic = slide.shapes.add_picture(image_stream, left, top, width=pic_width, height=pic_height)
                        top += pic.height + Inches(0.1)
                        log_func(f"Added LaTeX image for Option ক: {code}")
                elif '[option_খ]' in text:
                    option_text = row_data.get('Option_খ', '')
                    option_text_no_latex = remove_latex(option_text)
                    shape.text_frame.text = f"Option খ. {option_text_no_latex}"

                    # Apply font settings for Option_খ
                    apply_font_settings(shape, font_settings['Option'], log_func=log_func)

                    # Handle equations in option
                    option_latex_codes = extract_latex(option_text)
                    left = shape.left
                    top = shape.top + shape.height + Inches(0.1)
                    for code in option_latex_codes:
                        image = latex_to_image(code)
                        image_stream = io.BytesIO()
                        image.save(image_stream, format='PNG')
                        image_stream.seek(0)

                        # Add picture with desired size
                        pic_width = Inches(3)  # Adjust as needed
                        pic_height = Inches(0.8)  # Adjust as needed
                        pic = slide.shapes.add_picture(image_stream, left, top, width=pic_width, height=pic_height)
                        top += pic.height + Inches(0.1)
                        log_func(f"Added LaTeX image for Option খ: {code}")
                elif '[option_গ]' in text:
                    option_text = row_data.get('Option_গ', '')
                    option_text_no_latex = remove_latex(option_text)
                    shape.text_frame.text = f"Option গ. {option_text_no_latex}"

                    # Apply font settings for Option_গ
                    apply_font_settings(shape, font_settings['Option'], log_func=log_func)

                    # Handle equations in option
                    option_latex_codes = extract_latex(option_text)
                    left = shape.left
                    top = shape.top + shape.height + Inches(0.1)
                    for code in option_latex_codes:
                        image = latex_to_image(code)
                        image_stream = io.BytesIO()
                        image.save(image_stream, format='PNG')
                        image_stream.seek(0)

                        # Add picture with desired size
                        pic_width = Inches(3)  # Adjust as needed
                        pic_height = Inches(0.8)  # Adjust as needed
                        pic = slide.shapes.add_picture(image_stream, left, top, width=pic_width, height=pic_height)
                        top += pic.height + Inches(0.1)
                        log_func(f"Added LaTeX image for Option গ: {code}")
                elif '[option_ঘ]' in text:
                    option_text = row_data.get('Option_ঘ', '')
                    option_text_no_latex = remove_latex(option_text)
                    shape.text_frame.text = f"Option ঘ. {option_text_no_latex}"

                    # Apply font settings for Option_ঘ
                    apply_font_settings(shape, font_settings['Option'], log_func=log_func)

                    # Handle equations in option
                    option_latex_codes = extract_latex(option_text)
                    left = shape.left
                    top = shape.top + shape.height + Inches(0.1)
                    for code in option_latex_codes:
                        image = latex_to_image(code)
                        image_stream = io.BytesIO()
                        image.save(image_stream, format='PNG')
                        image_stream.seek(0)

                        # Add picture with desired size
                        pic_width = Inches(3)  # Adjust as needed
                        pic_height = Inches(0.8)  # Adjust as needed
                        pic = slide.shapes.add_picture(image_stream, left, top, width=pic_width, height=pic_height)
                        top += pic.height + Inches(0.1)
                        log_func(f"Added LaTeX image for Option ঘ: {code}")
                elif '[board_institute]' in text:
                    board_institute = row_data.get('Board/Institute', '')
                    shape.text_frame.text = f"{board_institute}"

                    # Apply font settings for Board/Institute
                    apply_font_settings(shape, font_settings['Board/Institute'], log_func=log_func)

                    log_func("Added Board/Institute to slide.")
                elif '[answer]' in text:
                    answer_text = row_data.get('Answer', '')
                    shape.text_frame.text = f"Answer: {answer_text}"

                    # Apply font settings for Answer
                    apply_font_settings(shape, font_settings['Answer'], log_func=log_func)

                    # Handle equations in answer (if any)
                    answer_latex_codes = extract_latex(answer_text)
                    left = shape.left
                    top = shape.top + shape.height + Inches(0.1)
                    for code in answer_latex_codes:
                        image = latex_to_image(code)
                        image_stream = io.BytesIO()
                        image.save(image_stream, format='PNG')
                        image_stream.seek(0)

                        # Add picture with desired size
                        pic_width = Inches(2)  # Adjust as needed
                        pic_height = Inches(0.5)  # Adjust as needed
                        pic = slide.shapes.add_picture(image_stream, left, top, width=pic_width, height=pic_height)
                        top += pic.height + Inches(0.1)
                        log_func(f"Added LaTeX image for Answer: {code}")

                    log_func("Added answer to slide.")

        # Optionally remove the template slide if you don't want it in the final presentation
        # prs.slides.remove(prs.slides[template_slide_index])

    # Save the presentation
    prs.save(output_pptx)
    log_func(f"PowerPoint presentation saved to {output_pptx}")


def apply_font_settings(shape, settings, log_func=None):
    """
    Apply font settings to all runs in all paragraphs of a shape's text frame.
    """
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            original_font = run.font.name
            try:
                run.font.name = settings['name'].strip()
            except Exception as e:
                run.font.name = 'Kalpurush'  # Fallback font
                if log_func:
                    log_func(f"Failed to set font '{settings['name']}'. Falling back to Kalpurush. Error: {e}")
            run.font.size = Pt(settings['size'])
            run.font.bold = settings['bold']
            run.font.italic = settings['italic']
            # Convert hex color code to RGBColor
            try:
                rgb = hex_to_rgb(settings['color'])
                run.font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
            except Exception as e:
                run.font.color.rgb = RGBColor(0, 0, 0)  # Default to black if conversion fails
                if log_func:
                    log_func(f"Error setting font color: {e}")
            if log_func:
                log_func(f"Applied font: '{run.font.name}' (was '{original_font}'), "
                         f"Size: {settings['size']}, Bold: {settings['bold']}, "
                         f"Italic: {settings['italic']}, Color: {settings['color']}")


def hex_to_rgb(hex_color):
    """
    Convert hex color string to RGB tuple.
    """
    hex_color = hex_color.lstrip('#')
    if len(hex_color) != 6:
        raise ValueError(f"Invalid hex color: {hex_color}")
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def latex_to_image(latex_code):
    """
    Render LaTeX code to an image and return it as a PIL Image object.
    """
    plt.rc('text', usetex=True)
    plt.rc('text.latex', preamble=r'\usepackage{amsmath}')
    fig = plt.figure(figsize=(0.01, 0.01))
    fig.text(0, 0, f'${latex_code}$', fontsize=12)
    buffer = io.BytesIO()
    plt.axis('off')
    plt.savefig(buffer, format='png', bbox_inches='tight', dpi=300, transparent=True)
    plt.close(fig)
    buffer.seek(0)
    image = Image.open(buffer)
    return image


def remove_latex(text):
    """
    Remove LaTeX code from text.
    """
    return re.sub(r'\$.*?\$', '', text)


def extract_latex(text):
    """
    Extract LaTeX code from text.
    """
    return re.findall(r'\$(.*?)\$', text)


def get_slide_layout_by_name(pres, name):
    """
    Retrieve a slide layout by its name.
    """
    for layout in pres.slide_layouts:
        if layout.name == name:
            return layout
    return None


def duplicate_slide(pres, index):
    """
    Duplicate a slide based on its index in the presentation.
    """
    source = pres.slides[index]
    blank_slide_layout = get_slide_layout_by_name(pres, 'Blank')
    if blank_slide_layout is None:
        # If 'Blank' layout is not found, use the first available layout
        blank_slide_layout = pres.slide_layouts[0]
    dest = pres.slides.add_slide(blank_slide_layout)
    for shape in source.shapes:
        # Copy each shape to the new slide
        el = shape.element
        new_el = copy.deepcopy(el)
        dest.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    return dest


def main():
    # This main function is kept empty as the GUI handles execution
    pass


if __name__ == '__main__':
    root = tk.Tk()
    app = PowerPointGeneratorGUI(root)
    root.mainloop()